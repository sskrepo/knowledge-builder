"""ArtifactComparator — text-only semantic comparison of produced vs reference artifacts.

Per ADR-029 (accepted 2026-05-15): text comparator only, NO vision-LLM.

Public interface (what Stream A's S5 imports):

    from framework.skill_builder.comparator import ArtifactComparator, ComparatorResult

    comparator = ArtifactComparator(llm=<LLMClient>)

    # Gate: detect image-only upload BEFORE calling compare()
    if comparator.is_image_only(artifact_bytes, artifact_type):
        # Hard-reject: surface ADR-029's prescribed message to the user.
        # This module returns the signal; the caller (S5 / conversation.py) routes.
        ...

    result: ComparatorResult = comparator.compare(
        reference_bytes=ref_bytes,
        produced_bytes=prod_bytes,
        artifact_type="pptx",  # "pptx" | "docx" | "md" | "txt"
    )

    # Result fields consumed by ADR-029 EVAL gap-report and failure-classifier:
    result.structure_score    # float 0.0–1.0: fraction of reference sections found
    result.density_score      # float 0.0–1.0: content volume ratio (word count)
    result.missing_sections   # list[str]: reference sections absent from produced
    result.thin_sections      # list[str]: sections present but density < 0.5x ref
    result.gap_report         # str: human-readable summary for the CHANGE PROPOSAL

Image-only detection contract (ADR-029 §C.1):
    is_image_only(artifact_bytes, artifact_type) -> bool
        Return True if the artifact has zero extractable text.
        Uses the same zero-text-shapes pattern as _analyze_pptx in analyze_artifact.py.
        Raises ValueError for unsupported artifact_type.

Synonym normalisation:
    Section-name matching uses a hardcoded synonym map (not an LLM call) to keep
    structure scoring deterministic.  The LLM is used ONLY for the optional
    semantic scoring path (compare()) when synonym-map misses are possible.

Design constraints:
    - NO dependency on conversation.py (standalone module).
    - NO vision-LLM path — all methods are text-only.
    - python-pptx is available (1.0.2+).  python-docx is optional; the module
      degrades gracefully with an ImportError warning.
    - Fail loud: every method raises ValueError with an actionable message rather
      than silently returning empty results.
"""
from __future__ import annotations

import io
import json
import logging
import re
from dataclasses import dataclass, field
from typing import Any

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Synonym map — deterministic section-name normalisation.
# Keys are canonical names (lowercase, stripped).  Values are lists of aliases.
# Extend this map as new synonym pairs are observed in the field.
# ---------------------------------------------------------------------------

_SYNONYM_GROUPS: list[list[str]] = [
    ["next steps", "action items", "actions", "follow-up", "follow up", "next actions"],
    ["key milestones", "milestones", "timeline", "schedule", "key dates", "dates"],
    ["risks", "risks & mitigations", "risks and mitigations", "risk register",
     "risk mitigations", "risks mitigations", "risks_mitigations"],
    ["status", "project status", "current status", "overall status", "rag status"],
    ["executive summary", "exec summary", "summary", "overview", "executive overview"],
    ["scope", "project scope", "in scope", "scope definition"],
    ["assumptions", "key assumptions", "assumptions & constraints", "constraints"],
    ["blockers", "blocking issues", "impediments", "issues"],
    ["budget", "financials", "cost", "spend", "cost summary"],
    ["team", "team members", "stakeholders", "resource plan"],
    ["objectives", "goals", "key objectives", "project goals", "okrs"],
    ["decisions", "decision log", "key decisions"],
    ["dependencies", "external dependencies", "key dependencies"],
    ["accomplishments", "completed", "achievements", "done", "completed this period"],
    ["orm status", "orm", "operational readiness", "go-live readiness"],
]

# Build a flat lookup: canonical_key -> canonical_name
# canonical_key = lowercase, stripped, underscores-collapsed alias
_SYNONYM_LOOKUP: dict[str, str] = {}
for _group in _SYNONYM_GROUPS:
    _canonical = _group[0]  # first entry is the canonical name
    for _alias in _group:
        _key = re.sub(r"[\s_\-]+", " ", _alias.lower()).strip()
        _SYNONYM_LOOKUP[_key] = _canonical


def _normalise_section(name: str) -> str:
    """Return the canonical name for a section, collapsing synonyms."""
    key = re.sub(r"[\s_\-]+", " ", name.lower()).strip()
    return _SYNONYM_LOOKUP.get(key, key)


# ---------------------------------------------------------------------------
# ComparatorResult — typed output consumed by ADR-029 EVAL + failure-classifier
# ---------------------------------------------------------------------------

@dataclass
class ComparatorResult:
    """Structured output from ArtifactComparator.compare().

    Per ADR-029 §C.2 rubric:
        structure_score: fraction of reference sections found in produced artifact.
        density_score:   average word-count ratio (produced/reference) per matching
                         section; 1.0 means equal density; >1.0 means produced is
                         denser (not penalised).
        missing_sections: section names (as they appear in reference) that are
                         entirely absent from the produced artifact.
        thin_sections:   section names present in produced but whose word-count is
                         less than 0.5x the corresponding reference section.
        gap_report:      human-readable summary for surfacing in the CHANGE PROPOSAL
                         turn (must_show_human=True).  Example:
                         "The produced PPTX has 5 sections; your reference had 7.
                          Missing: Risks, Next Steps.  Thin sections: Status."
    """
    structure_score: float
    density_score: float
    missing_sections: list[str] = field(default_factory=list)
    thin_sections: list[str] = field(default_factory=list)
    gap_report: str = ""

    def to_dict(self) -> dict[str, Any]:
        return {
            "structure_score": self.structure_score,
            "density_score": self.density_score,
            "missing_sections": self.missing_sections,
            "thin_sections": self.thin_sections,
            "gap_report": self.gap_report,
        }


# ---------------------------------------------------------------------------
# Image-only / unsupported result signal
# ---------------------------------------------------------------------------

# ADR-029 §C.1 hard-reject message — callers surface this verbatim.
IMAGE_ONLY_MESSAGE = (
    "Image-based reference artifacts are not supported yet (no Vision-LLM backend). "
    "Please upload a text-bearing reference (text-extractable PPTX/DOCX/MD). "
    "The current reference will be discarded."
)

SUPPORTED_TYPES = frozenset({"pptx", "docx", "md", "txt"})


# ---------------------------------------------------------------------------
# Text extraction helpers — mirroring analyze_artifact.py patterns
# ---------------------------------------------------------------------------

def _extract_pptx_sections(artifact_bytes: bytes) -> list[tuple[str, int]]:
    """Extract (title, word_count) pairs from a PPTX byte stream.

    Returns a list of (raw_slide_title, word_count_of_all_text_on_slide).
    Slides with no title are skipped for structure scoring (but their text
    is counted against overall density).

    Raises ImportError if python-pptx is not installed (should not happen
    in this repo — python-pptx 1.0.2 is a known dependency).
    """
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(io.BytesIO(artifact_bytes))
    sections: list[tuple[str, int]] = []
    for slide in prs.slides:
        title_text = ""
        if slide.shapes.title and slide.shapes.title.text:
            title_text = slide.shapes.title.text.strip()

        # Count all words across all text frames on this slide
        words = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                words += len(shape.text_frame.text.split())

        if title_text:
            sections.append((title_text, words))

    return sections


def _count_pptx_text_shapes(artifact_bytes: bytes) -> int:
    """Return total text-shape count across all slides — mirrors _analyze_pptx.

    Zero means image-only.  This is the gate for the ADR-029 hard-reject path.
    """
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(io.BytesIO(artifact_bytes))
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                total += 1
    return total


def _extract_docx_sections(artifact_bytes: bytes) -> list[tuple[str, int]]:
    """Extract (heading, word_count) pairs from a DOCX byte stream.

    Raises ImportError if python-docx is not installed.
    """
    try:
        from docx import Document  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-docx is not installed; DOCX comparison is unavailable. "
            "Install it with: pip install python-docx"
        ) from exc

    doc = Document(io.BytesIO(artifact_bytes))
    sections: list[tuple[str, int]] = []
    current_heading: str | None = None
    current_words = 0

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            if current_heading is not None:
                sections.append((current_heading, current_words))
            current_heading = para.text.strip()
            current_words = 0
        else:
            current_words += len(para.text.split())

    if current_heading is not None:
        sections.append((current_heading, current_words))

    return sections


def _count_docx_text(artifact_bytes: bytes) -> int:
    """Return total word count of a DOCX — used for image-only detection."""
    try:
        from docx import Document  # type: ignore[import]
    except ImportError:
        log.warning("python-docx not installed; DOCX text count unavailable")
        return 1  # assume non-empty when we cannot check

    doc = Document(io.BytesIO(artifact_bytes))
    total = 0
    for para in doc.paragraphs:
        total += len(para.text.split())
    return total


def _extract_md_sections(artifact_bytes: bytes) -> list[tuple[str, int]]:
    """Extract (heading, word_count) pairs from a Markdown byte stream."""
    text = artifact_bytes.decode("utf-8", errors="replace")
    lines = text.splitlines()
    sections: list[tuple[str, int]] = []
    current_heading: str | None = None
    current_words = 0

    for line in lines:
        m = re.match(r"^#{1,6}\s+(.+)$", line)
        if m:
            if current_heading is not None:
                sections.append((current_heading, current_words))
            current_heading = m.group(1).strip()
            current_words = 0
        elif current_heading is not None:
            current_words += len(line.split())

    if current_heading is not None:
        sections.append((current_heading, current_words))

    return sections


def _count_md_text(artifact_bytes: bytes) -> int:
    """Return total word count of a Markdown file."""
    text = artifact_bytes.decode("utf-8", errors="replace")
    return len(text.split())


def _extract_sections(artifact_bytes: bytes, artifact_type: str) -> list[tuple[str, int]]:
    """Dispatch to the right section extractor based on artifact_type.

    Returns list of (section_title, word_count) tuples.
    Raises ValueError for unsupported types.
    """
    t = artifact_type.lower().lstrip(".")
    if t == "pptx":
        return _extract_pptx_sections(artifact_bytes)
    if t == "docx":
        return _extract_docx_sections(artifact_bytes)
    if t in ("md", "txt"):
        return _extract_md_sections(artifact_bytes)
    raise ValueError(
        f"ArtifactComparator: unsupported artifact_type={artifact_type!r}. "
        f"Supported types: {sorted(SUPPORTED_TYPES)}. "
        "If this is a new format, add a section extractor in comparator.py."
    )


# ---------------------------------------------------------------------------
# ArtifactComparator
# ---------------------------------------------------------------------------

class ArtifactComparator:
    """Semantic text comparator for produced vs reference artifacts.

    Per ADR-029: text-only, no vision-LLM.  Structure and density scores are
    computed deterministically (word-count + synonym-map).  The optional LLM
    parameter is reserved for future semantic scoring extensions; all current
    scoring is deterministic.

    Usage:
        comparator = ArtifactComparator(llm=llm_client)  # llm may be None

        # 1. Gate: detect image-only before calling compare()
        if comparator.is_image_only(ref_bytes, "pptx"):
            # surface IMAGE_ONLY_MESSAGE to the user — do not advance state

        # 2. Score
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
    """

    def __init__(self, llm=None) -> None:
        """
        Args:
            llm: optional LLM client (framework LLMClient or compatible).
                 Currently unused in the deterministic scoring path.
                 Reserved for future semantic-scoring extensions (e.g. synonym
                 expansion beyond the hardcoded map).
        """
        self._llm = llm

    # ------------------------------------------------------------------
    # is_image_only — ADR-029 hard-reject gate
    # ------------------------------------------------------------------

    def is_image_only(self, artifact_bytes: bytes, artifact_type: str) -> bool:
        """Return True if the artifact contains zero extractable text.

        Uses the same zero-text-shapes pattern as analyze_artifact._analyze_pptx:
        count text-bearing shapes across all slides; return True if count == 0.

        For DOCX / MD / TXT: return True if total word count == 0.

        Args:
            artifact_bytes: raw bytes of the uploaded artifact.
            artifact_type:  "pptx" | "docx" | "md" | "txt" (case-insensitive).

        Returns:
            True  — artifact is image-only / has no extractable text.
            False — artifact has text and can be compared.

        Raises:
            ValueError: for unsupported artifact_type.
            ImportError: if python-pptx is not installed (should not happen in
                         this repo — it is a declared dependency).
        """
        t = artifact_type.lower().lstrip(".")
        if t not in SUPPORTED_TYPES:
            raise ValueError(
                f"ArtifactComparator.is_image_only: unsupported artifact_type={artifact_type!r}. "
                f"Supported: {sorted(SUPPORTED_TYPES)}."
            )

        if t == "pptx":
            try:
                count = _count_pptx_text_shapes(artifact_bytes)
                return count == 0
            except ImportError as exc:
                raise ImportError(
                    "python-pptx is required for PPTX comparison. "
                    "It is listed in framework/requirements.txt."
                ) from exc

        if t == "docx":
            try:
                count = _count_docx_text(artifact_bytes)
                return count == 0
            except ImportError:
                # python-docx not installed: cannot determine; log a warning and
                # treat as non-image-only (conservative — do not hard-reject when
                # we cannot check).
                log.warning(
                    "is_image_only: python-docx not installed; treating DOCX as "
                    "non-image-only (conservative default).  Install python-docx "
                    "for accurate DOCX support."
                )
                return False

        if t in ("md", "txt"):
            count = _count_md_text(artifact_bytes)
            return count == 0

        # Should not reach here (validated above) but be explicit.
        raise ValueError(f"Unhandled artifact_type: {artifact_type!r}")

    # ------------------------------------------------------------------
    # compare — main scoring entry point
    # ------------------------------------------------------------------

    def compare(
        self,
        reference_bytes: bytes,
        produced_bytes: bytes,
        artifact_type: str,
    ) -> ComparatorResult:
        """Compare produced artifact against reference and return a scored result.

        Scoring rubric (ADR-029 §C.2):
            structure_score: fraction of reference sections present in produced.
                             Synonym normalisation applied before matching.
            density_score:   average word-count ratio (produced/reference) per
                             matched section, capped at 1.0.  Unmatched reference
                             sections contribute 0 to the average.
            missing_sections: reference sections absent from produced.
            thin_sections:   sections whose density ratio < 0.5.
            gap_report:      human-readable summary for the CHANGE PROPOSAL turn.

        Args:
            reference_bytes: bytes of the reference (uploaded by user).
            produced_bytes:  bytes of the produced artifact (from workflow run).
            artifact_type:   "pptx" | "docx" | "md" | "txt" (case-insensitive).

        Returns:
            ComparatorResult

        Raises:
            ValueError: if artifact_type is unsupported, or if the reference has
                        no sections (empty artifact — callers should call
                        is_image_only() first to handle that case).
        """
        t = artifact_type.lower().lstrip(".")
        if t not in SUPPORTED_TYPES:
            raise ValueError(
                f"ArtifactComparator.compare: unsupported artifact_type={artifact_type!r}. "
                f"Supported: {sorted(SUPPORTED_TYPES)}."
            )

        ref_sections = _extract_sections(reference_bytes, t)
        prod_sections = _extract_sections(produced_bytes, t)

        if not ref_sections:
            raise ValueError(
                "ArtifactComparator.compare: reference artifact has no extractable "
                "sections (empty or image-only).  Call is_image_only() before compare() "
                "and handle the image-only case at the call site."
            )

        # Build lookup: canonical_name -> word_count for produced artifact
        prod_lookup: dict[str, int] = {}
        for title, wc in prod_sections:
            canon = _normalise_section(title)
            # If multiple slides map to the same canonical name, sum their words
            prod_lookup[canon] = prod_lookup.get(canon, 0) + wc

        missing_sections: list[str] = []
        thin_sections: list[str] = []
        density_ratios: list[float] = []

        for ref_title, ref_wc in ref_sections:
            canon = _normalise_section(ref_title)
            if canon not in prod_lookup:
                missing_sections.append(ref_title)
                density_ratios.append(0.0)
            else:
                prod_wc = prod_lookup[canon]
                # Avoid division by zero for zero-word reference sections
                if ref_wc > 0:
                    ratio = min(prod_wc / ref_wc, 1.0)
                else:
                    ratio = 1.0  # reference section had no words — treat as matched
                density_ratios.append(ratio)
                if ratio < 0.5:
                    thin_sections.append(ref_title)

        n_ref = len(ref_sections)
        n_found = n_ref - len(missing_sections)
        structure_score = round(n_found / n_ref, 4) if n_ref > 0 else 1.0
        density_score = round(sum(density_ratios) / len(density_ratios), 4) if density_ratios else 1.0

        gap_report = self._build_gap_report(
            artifact_type=t,
            n_ref=n_ref,
            n_prod=len(prod_sections),
            missing_sections=missing_sections,
            thin_sections=thin_sections,
            structure_score=structure_score,
            density_score=density_score,
        )

        return ComparatorResult(
            structure_score=structure_score,
            density_score=density_score,
            missing_sections=missing_sections,
            thin_sections=thin_sections,
            gap_report=gap_report,
        )

    # ------------------------------------------------------------------
    # gap_report builder
    # ------------------------------------------------------------------

    def _build_gap_report(
        self,
        artifact_type: str,
        n_ref: int,
        n_prod: int,
        missing_sections: list[str],
        thin_sections: list[str],
        structure_score: float,
        density_score: float,
    ) -> str:
        """Build a human-readable gap report for the CHANGE PROPOSAL turn.

        The report is surfaced to the user with must_show_human=True.
        Keep it concise and actionable — persona teams are not engineers.
        """
        kind_label = {
            "pptx": "PPTX",
            "docx": "DOCX",
            "md": "document",
            "txt": "document",
        }.get(artifact_type, "artifact")

        lines: list[str] = []
        lines.append(
            f"The produced {kind_label} has {n_prod} section(s); "
            f"your reference had {n_ref}."
        )
        lines.append(
            f"Structure score: {structure_score:.0%}. "
            f"Content density score: {density_score:.0%}."
        )

        if missing_sections:
            joined = ", ".join(missing_sections)
            lines.append(f"Missing: {joined}.")

        if thin_sections:
            joined = ", ".join(thin_sections)
            lines.append(f"Thin sections (less than 50% of reference content): {joined}.")

        if not missing_sections and not thin_sections:
            lines.append(
                "No structural gaps detected. "
                "Review the content quality and field values manually."
            )

        return "  ".join(lines)
