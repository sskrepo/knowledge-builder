"""analyze_artifact — parse PPT/DOCX/email and infer required fields.

Per ADR-015 + ADR-007 amend 4.

Updated for Phase 3: now returns tuple[list[str], dict | None].
- list[str]: inferred field names
- dict | None: slide_mapping / section_mapping derived from artifact structure
  (slide titles for PPTX, heading structure for DOCX, section headings for MD).
  None when no structural mapping could be derived.
"""
from __future__ import annotations

import logging
import re
from pathlib import Path

log = logging.getLogger(__name__)


def analyze_artifact(path: str) -> tuple[list[str], dict | None]:
    """Parse an artifact and return (fields, mapping).

    fields: list of snake_case field name strings inferred from the artifact.
    mapping: dict mapping each field to its source location in the artifact
             (slide number, heading name, etc.) or None if not determinable.
    """
    p = Path(path)
    if not p.exists():
        log.warning("artifact path does not exist: %s; falling back to heuristic", path)
        return ["title", "summary", "details"], None

    if p.suffix == ".pptx":
        return _analyze_pptx(p)
    if p.suffix == ".docx":
        return _analyze_docx(p)
    if p.suffix in (".md", ".txt"):
        return _analyze_markdown(p)
    return ["title", "summary", "details"], None


def _analyze_pptx(p: Path) -> tuple[list[str], dict | None]:
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python-pptx not installed; using fallback fields")
        fields = ["title", "section_a", "section_b", "section_c"]
        return fields, None

    prs = Presentation(str(p))
    fields: list[str] = ["title"]
    mapping: dict = {
        "title": {"kind": "slide_title", "slide": 0, "raw_title": "Title Slide"},
    }

    total_text_shapes = 0
    for i, slide in enumerate(prs.slides):
        raw_title = ""
        if slide.shapes.title and slide.shapes.title.text:
            raw_title = slide.shapes.title.text.strip()
        # Scan all shapes for any text, even if no slide title is present
        slide_text_shapes = sum(
            1 for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()
        )
        total_text_shapes += slide_text_shapes

        if not raw_title:
            continue
        field = _to_field_name(raw_title)
        if field and field not in fields:
            fields.append(field)
            # Capture body text (all non-title text frames) for LLM description synthesis
            body_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        body_parts.append(t)
            body_text = " ".join(body_parts)[:400]
            mapping[field] = {
                "kind": "slide_title",
                "slide": i,
                "raw_title": raw_title,
                "body_text": body_text,
            }

    # ADR-026 Fix 1: image-only PPTX (no text shapes found) produces no usable fields.
    # Hard-fail with a clear actionable error instead of silently returning a
    # nearly-empty field list that triggers keyword heuristics downstream.
    # Vision-LLM analysis for image-only slides is deferred to ADR-027.
    if total_text_shapes == 0:
        raise ValueError(
            f"The uploaded PPTX '{p.name}' contains no text shapes "
            f"({len(prs.slides)} slide(s), all image-only). "
            "The framework cannot extract field names from an image-only presentation. "
            "Options:\n"
            "  1. Supply a text-bearing PPTX, DOCX, or Markdown file that represents "
            "the output structure.\n"
            "  2. Enter field names manually as a comma-separated list "
            "(e.g. 'scope, assumptions, status_bullets, next_steps, key_milestones, "
            "risks_mitigations, orm_status').\n"
            "Vision-LLM analysis of image slides is planned for ADR-027."
        )

    return fields, mapping if len(mapping) > 1 else None


def _analyze_docx(p: Path) -> tuple[list[str], dict | None]:
    try:
        from docx import Document
    except ImportError:
        return ["title", "introduction", "details", "conclusion"], None

    doc = Document(str(p))
    fields: list[str] = ["title"]
    mapping: dict = {
        "title": {"kind": "document_title", "heading_level": 0, "raw_heading": "Document Title"},
    }

    paragraphs = list(doc.paragraphs)
    for idx, para in enumerate(paragraphs):
        if not para.style.name.startswith("Heading"):
            continue
        raw_heading = para.text.strip()
        if not raw_heading:
            continue
        level_match = re.search(r"(\d+)", para.style.name)
        level = int(level_match.group(1)) if level_match else 1
        field = _to_field_name(raw_heading)
        if field and field not in fields:
            fields.append(field)
            # Capture first non-heading paragraphs under this heading as body_text
            body_parts: list[str] = []
            for body_para in paragraphs[idx + 1:]:
                if body_para.style.name.startswith("Heading"):
                    break
                t = body_para.text.strip()
                if t:
                    body_parts.append(t)
                if sum(len(p) for p in body_parts) > 400:
                    break
            body_text = " ".join(body_parts)[:400]
            mapping[field] = {
                "kind": "heading",
                "heading_level": level,
                "raw_heading": raw_heading,
                "body_text": body_text,
            }

    return fields, mapping if len(mapping) > 1 else None


def _analyze_markdown(p: Path) -> tuple[list[str], dict | None]:
    fields: list[str] = []
    mapping: dict = {}
    section_index = 0

    all_lines = p.read_text().splitlines()
    for i, line in enumerate(all_lines):
        m = re.match(r"^(#{1,6})\s+(.+)$", line)
        if not m:
            continue
        level = len(m.group(1))
        raw_heading = m.group(2).strip()
        field = _to_field_name(raw_heading)
        if field and field not in fields:
            fields.append(field)
            # Capture non-heading lines immediately following as body_text
            body_parts: list[str] = []
            for body_line in all_lines[i + 1:]:
                if re.match(r"^#{1,6}\s", body_line):
                    break
                t = body_line.strip()
                if t:
                    body_parts.append(t)
                if sum(len(x) for x in body_parts) > 400:
                    break
            body_text = " ".join(body_parts)[:400]
            mapping[field] = {
                "kind": "heading",
                "heading_level": level,
                "line_number": i + 1,
                "raw_heading": raw_heading,
                "section_index": section_index,
                "body_text": body_text,
            }
            section_index += 1

    if not fields:
        return ["title", "summary"], None

    return fields, mapping or None


def _to_field_name(text: str) -> str:
    t = text.lower().replace(" ", "_").replace("-", "_")
    t = re.sub(r"[^a-z0-9_]", "", t)
    t = re.sub(r"_+", "_", t).strip("_")
    return t[:60]
