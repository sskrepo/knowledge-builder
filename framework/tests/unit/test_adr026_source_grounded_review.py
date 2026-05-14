"""Unit tests for ADR-026 fixes.

Coverage:
  Fix 1 — analyze_artifact raises ValueError for image-only PPTX
  Fix 2 — sampler.fetch_samples calls live Confluence when page_id present
  Fix 3 — review_extractions calls LLM; raises RuntimeError when llm=None
  Fix 4 — _source_grounded_review fires at REVIEW_SCHEMA when sources are configured
  Fix 5 — PptxRenderer routes to weekly_exec_review_v1 layout; builds single slide
"""
from __future__ import annotations

import io
import json
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# Fix 1 — analyze_artifact hard-fail for image-only PPTX
# ---------------------------------------------------------------------------


def _make_image_only_pptx(tmp_path: Path) -> Path:
    """Build a minimal PPTX with a single PICTURE shape, no text."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    # Add a picture placeholder indirectly via add_shape (no text)
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(4))
    out = tmp_path / "image_only.pptx"
    prs.save(str(out))
    return out


class TestAnalyzeArtifactImageOnly:
    def test_image_only_pptx_raises_value_error(self, tmp_path):
        pptx_path = _make_image_only_pptx(tmp_path)
        from framework.skill_builder.analyze_artifact import analyze_artifact

        with pytest.raises(ValueError, match="image-only"):
            analyze_artifact(str(pptx_path))

    def test_text_pptx_does_not_raise(self, tmp_path):
        """A PPTX with at least one text shape should succeed."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Project Status"
        out = tmp_path / "with_text.pptx"
        prs.save(str(out))

        from framework.skill_builder.analyze_artifact import analyze_artifact

        fields, mapping = analyze_artifact(str(out))
        assert "project_status" in fields

    def test_missing_file_returns_fallback_fields(self):
        """Non-existent file returns fallback without raising."""
        from framework.skill_builder.analyze_artifact import analyze_artifact

        fields, mapping = analyze_artifact("/tmp/nonexistent_adr026_test.pptx")
        assert "title" in fields


# ---------------------------------------------------------------------------
# Fix 2 — sampler live Confluence fetch
# ---------------------------------------------------------------------------


class TestSamplerLiveConfluenceFetch:
    def test_live_fetch_called_when_page_id_present(self, tmp_path):
        """fetch_samples calls _fetch_confluence_live when page_id given."""
        from framework.skill_builder import sampler

        mock_sample = {
            "source_citation": "https://confluence.example.com?pageId=12345",
            "title": "26ai Status",
            "content": "## FA DB Upgrade\n\nScope: upgrade Oracle FA DB to 26ai.",
            "_live": True,
        }

        with patch.object(sampler, "_fetch_confluence_live", return_value=[mock_sample]) as mock_live:
            result = sampler.fetch_samples(
                adapter_name="confluence",
                source_query={"page_id": "12345"},
                n=1,
                kbf_env="laptop",
                repo_root=tmp_path,
            )

        mock_live.assert_called_once()
        assert result == [mock_sample]

    def test_fallback_to_fixtures_when_live_fails(self, tmp_path):
        """Falls back to fixtures when live fetch raises, not require_live."""
        from framework.skill_builder import sampler

        fixture_dir = tmp_path / "_dev_fixtures" / "confluence"
        fixture_dir.mkdir(parents=True)
        (fixture_dir / "sample.json").write_text(
            json.dumps({"source_citation": "fixture://confluence/sample.json", "content": "hello"})
        )

        # Patch FIXTURES_ROOT to tmp_path/_dev_fixtures
        with patch.object(sampler, "FIXTURES_ROOT", tmp_path / "_dev_fixtures"):
            with patch.object(sampler, "_fetch_confluence_live", side_effect=RuntimeError("no adapter")):
                result = sampler.fetch_samples(
                    adapter_name="confluence",
                    source_query={"page_id": "99999"},
                    n=1,
                    kbf_env="laptop",
                    repo_root=tmp_path,
                )

        assert result[0]["source_citation"] == "fixture://confluence/sample.json"

    def test_require_live_raises_when_adapter_unavailable(self, tmp_path):
        """require_live=True propagates the error."""
        from framework.skill_builder import sampler

        with patch.object(sampler, "_fetch_confluence_live", side_effect=RuntimeError("no adapter")):
            with pytest.raises(RuntimeError, match="live fetch required"):
                sampler.fetch_samples(
                    adapter_name="confluence",
                    source_query={"page_id": "12345"},
                    n=1,
                    require_live=True,
                    kbf_env="laptop",
                    repo_root=tmp_path,
                )

    def test_no_page_id_uses_fixtures(self, tmp_path, monkeypatch):
        """Without page_id/page_url, fixtures are used (no live fetch)."""
        from framework.skill_builder import sampler

        monkeypatch.setenv("KBF_STORE_BACKEND", "filestore")
        fixture_dir = tmp_path / "_dev_fixtures" / "confluence"
        fixture_dir.mkdir(parents=True)
        (fixture_dir / "data.json").write_text(
            json.dumps({"source_citation": "fixture://confluence/data.json", "content": "text"})
        )
        with patch.object(sampler, "FIXTURES_ROOT", tmp_path / "_dev_fixtures"):
            with patch.object(sampler, "_fetch_confluence_live") as mock_live:
                result = sampler.fetch_samples(
                    adapter_name="confluence",
                    source_query={"space": "FAAAS"},
                    n=1,
                    kbf_env="laptop",
                    repo_root=tmp_path,
                )
        mock_live.assert_not_called()


# ---------------------------------------------------------------------------
# Fix 3 — review_extractions LLM path + hard-fail
# ---------------------------------------------------------------------------


class TestReviewExtractions:
    def _make_schema(self) -> dict:
        return {
            "properties": {
                "project_name": {"type": "string", "description": "Project name"},
                "status_bullets": {"type": "array", "description": "Status bullets"},
            },
            "required": ["project_name"],
        }

    def _make_samples(self) -> list[dict]:
        return [
            {
                "source_citation": "https://confluence.example.com?pageId=123",
                "content": "FA DB Upgrade to 26ai — Status: In Progress",
            }
        ]

    def test_raises_without_llm(self):
        from framework.skill_builder.review import review_extractions

        with pytest.raises(RuntimeError, match="llm is required"):
            review_extractions(self._make_samples(), self._make_schema(), llm=None)

    def test_stub_mode_works_without_llm(self):
        from framework.skill_builder.review import review_extractions

        result = review_extractions(
            self._make_samples(), self._make_schema(), llm=None, stub_mode=True
        )
        assert "extractions" in result
        assert result["extraction_mode"] == "stub"

    def test_llm_path_called_when_wired(self):
        from framework.skill_builder.review import review_extractions

        mock_llm = MagicMock()
        mock_llm.chat.return_value = {
            "text": json.dumps({"project_name": "FA DB 26ai", "status_bullets": ["In Progress"]})
        }
        result = review_extractions(self._make_samples(), self._make_schema(), llm=mock_llm)

        assert result["extraction_mode"] == "llm"
        assert mock_llm.chat.called
        extractions = result["extractions"]
        assert len(extractions) == 1
        assert extractions[0]["extracted"]["project_name"] == "FA DB 26ai"

    def test_llm_path_reports_missing_fields(self):
        from framework.skill_builder.review import review_extractions

        mock_llm = MagicMock()
        # LLM returns empty object — project_name is required but missing
        mock_llm.chat.return_value = {"text": json.dumps({})}
        result = review_extractions(self._make_samples(), self._make_schema(), llm=mock_llm)

        issues = result["issues"]
        assert any("project_name" in issue for issue in issues)

    def test_field_coverage_calculated_correctly(self):
        from framework.skill_builder.review import review_extractions

        mock_llm = MagicMock()
        mock_llm.chat.return_value = {
            "text": json.dumps({"project_name": "26ai", "status_bullets": ["done"]})
        }
        result = review_extractions(self._make_samples(), self._make_schema(), llm=mock_llm)
        assert result["field_coverage"]["project_name"] == 1.0


# ---------------------------------------------------------------------------
# Fix 4 — _source_grounded_review in conversation
# ---------------------------------------------------------------------------


class TestSourceGroundedReview:
    def _make_conversation_at_review_schema(self):
        from framework.skill_builder.conversation import SkillBuilderConversation

        skill_store = MagicMock()
        llm = MagicMock()
        conv = SkillBuilderConversation(
            persona="tpm",
            user_id="test-user",
            skill_store=skill_store,
            llm=llm,
        )
        conv._data.intent_description = "weekly exec review for 26ai project"
        conv._data.skill_name = "weekly_exec_review_26ai"
        conv._data.fields = ["scope", "status_bullets", "key_milestones"]
        conv._data.field_specs = {
            "scope": {"type": "string", "description": "Project scope"},
            "status_bullets": {"type": "array", "description": "Status updates"},
            "key_milestones": {"type": "array", "description": "Key milestones"},
        }
        conv._data.sources = [
            {
                "kind": "confluence",
                "pages": ["https://confluence.example.com/pages/viewpage.action?pageId=20030556732"],
            }
        ]
        return conv

    def test_source_review_fires_and_is_in_turn_data(self):
        """_source_grounded_review result appears in the REVIEW_SCHEMA turn data."""
        conv = self._make_conversation_at_review_schema()

        # Mock sampler to return a live sample
        live_sample = {
            "source_citation": "https://confluence.example.com?pageId=20030556732",
            "content": "## FA DB Upgrade to 26ai\n\nScope: Upgrade FA DB from 19c to 26ai.\n\nKey Milestones: ...",
            "title": "FA DB Upgrade 26ai",
            "_live": True,
        }
        review_result = {
            "unsupportable_fields": [],
            "suggested_additions": [
                {"field": "orm_status", "type": "string", "description": "ORM status",
                 "reason": "Source mentions ORM approval status"}
            ],
            "enum_corrections": [],
            "summary": "Schema is mostly aligned; consider adding orm_status.",
        }

        with patch("framework.skill_builder.conversation.fetch_samples", return_value=[live_sample]):
            with patch.object(conv._llm, "chat") as mock_chat:
                mock_chat.return_value = {"text": json.dumps(review_result)}
                turn = conv._advance_to_review_schema()

        assert turn.state == "REVIEW_SCHEMA"
        assert turn.data is not None
        sr = turn.data.get("source_review")
        assert sr is not None, "source_review should be in turn.data"
        assert sr["summary"] == review_result["summary"]
        assert len(sr["suggested_additions"]) == 1
        assert "orm_status" in turn.message

    def test_source_review_skipped_when_no_page_sources(self):
        """No source review when sources have no page IDs."""
        conv = self._make_conversation_at_review_schema()
        conv._data.sources = [{"kind": "confluence", "space": "FAAAS"}]

        with patch("framework.skill_builder.conversation.fetch_samples") as mock_sampler:
            turn = conv._advance_to_review_schema()

        mock_sampler.assert_not_called()
        assert turn.state == "REVIEW_SCHEMA"
        # source_review should be None or absent
        sr = (turn.data or {}).get("source_review")
        assert sr is None

    def test_source_review_skipped_when_no_llm(self):
        """No source review when LLM is not wired."""
        from framework.skill_builder.conversation import SkillBuilderConversation

        skill_store = MagicMock()
        conv = SkillBuilderConversation(
            persona="tpm",
            user_id="test-user",
            skill_store=skill_store,
            llm=None,  # no LLM
        )
        conv._data.intent_description = "weekly review"
        conv._data.skill_name = "test_skill"
        conv._data.fields = ["scope"]
        conv._data.sources = [{"kind": "confluence", "pages": ["https://example.com?pageId=123"]}]

        with patch("framework.skill_builder.conversation.fetch_samples") as mock_sampler:
            turn = conv._advance_to_review_schema()

        mock_sampler.assert_not_called()

    def test_source_review_failure_is_non_blocking(self):
        """If source-grounded review raises, REVIEW_SCHEMA still renders."""
        conv = self._make_conversation_at_review_schema()

        with patch("framework.skill_builder.conversation.fetch_samples",
                   side_effect=RuntimeError("Confluence down")):
            turn = conv._advance_to_review_schema()

        assert turn.state == "REVIEW_SCHEMA"
        # source_review should be absent (None) — failure is advisory
        sr = (turn.data or {}).get("source_review")
        assert sr is None


# ---------------------------------------------------------------------------
# Fix 5 — PptxRenderer layout-aware dispatch
# ---------------------------------------------------------------------------


class TestPptxRendererLayoutAware:
    def _make_26ai_data(self) -> dict:
        return {
            "title": "FA DB Upgrade to 26ai",
            "layout": "weekly_exec_review_v1",
            "jira_id": "FAAASPMO-1190",
            "scope": "Upgrade Oracle FA DB from 19c to 26ai in production.",
            "assumptions": ["Migration window approved", "DBA team available"],
            "status_bullets": [
                "Completed: Dev environment upgraded",
                "In Progress: UAT testing",
                "On Hold: Prod cutover pending sign-off",
            ],
            "next_steps": ["Complete UAT by 2026-05-20", "Schedule prod window"],
            "key_milestones": [
                "Exa Infra plan baseline — 2026-05-20",
                "UAT complete — 2026-05-28 (At Risk)",
                "Prod cutover — 2026-06-10",
            ],
            "orm_status": "ORM FAAASPMO-1190 submitted, approval pending.",
            "risks_mitigations": [
                "Risk: UAT timeline; Mitigation: daily standups + escalation path",
            ],
            "sections": {},
            "extracted": {
                "jira_id": "FAAASPMO-1190",
                "project_name": "FA DB Upgrade to 26ai",
            },
            "citations": ["https://confluence.oraclecorp.com/pages/viewpage.action?pageId=20030556732"],
            "generated_at": "2026-05-13T12:00:00Z",
        }

    def test_layout_weekly_exec_review_v1_produces_one_slide(self):
        """weekly_exec_review_v1 layout builds exactly one slide."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        from framework.renderers.pptx_renderer import PptxRenderer

        renderer = PptxRenderer()
        pptx_bytes = renderer.render(self._make_26ai_data())

        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 1, f"Expected 1 slide, got {len(prs.slides)}"

    def test_layout_weekly_exec_review_v1_has_oracle_header(self):
        """The slide should have the Oracle red header band shape."""
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
        except ImportError:
            pytest.skip("python-pptx not installed")

        from framework.renderers.pptx_renderer import PptxRenderer

        renderer = PptxRenderer()
        pptx_bytes = renderer.render(self._make_26ai_data())

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide = prs.slides[0]
        # Check for the red header shape
        red_shapes = [
            s for s in slide.shapes
            if hasattr(s, "fill") and s.fill.type is not None
            and s.fill.type.name == "SOLID"
            and s.fill.fore_color.rgb == RGBColor(199, 70, 52)
        ]
        assert len(red_shapes) >= 1, "Expected at least one Oracle red shape"

    def test_default_layout_produces_multiple_slides(self):
        """Without layout key, the default multi-slide path is taken."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        from framework.renderers.pptx_renderer import PptxRenderer

        data = {
            "title": "Test",
            "sections": {"Scope": "text", "Status": "running", "Next Steps": "deploy"},
        }
        renderer = PptxRenderer()
        pptx_bytes = renderer.render(data)

        prs = Presentation(io.BytesIO(pptx_bytes))
        # Default path: title slide + 3 section slides = 4
        assert len(prs.slides) >= 2

    def test_layout_aware_contains_jira_id_text(self):
        """Jira ID should appear in the slide text somewhere."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        from framework.renderers.pptx_renderer import PptxRenderer

        renderer = PptxRenderer()
        pptx_bytes = renderer.render(self._make_26ai_data())

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide = prs.slides[0]
        all_text = " ".join(
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "FAAASPMO-1190" in all_text, f"Jira ID not found in slide text: {all_text[:300]}"
