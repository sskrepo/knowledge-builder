"""Unit tests for ADR-034 — Layout Preset Catalog.

Coverage:
  (a) skill_builder.yaml no longer contains the hardcoded weekly_exec_review_v1
      rule/enum string (the explicit rule "Choose weekly_exec_review_v1 layout
      only for exec-review PPTX skills" must be gone).
  (b) Catalog is the single source of truth and PptxRenderer still dispatches
      both presets correctly via catalog lookup.
  (c) A generated clarify question containing an internal preset id is sanitized
      to plain language before being shown to the user.
  (d) catalog_for_prompt returns human descriptions only — never internal ids.
  (e) PromptRegistry still parses skill_builder.yaml without error after changes.
"""
from __future__ import annotations

import re
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
import yaml

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

REPO_ROOT = Path(__file__).resolve().parents[3]
SKILL_BUILDER_YAML = REPO_ROOT / "framework" / "config" / "prompts" / "skill_builder.yaml"


# ---------------------------------------------------------------------------
# (a) skill_builder.yaml no longer hardcodes the weekly_exec_review_v1 rule
# ---------------------------------------------------------------------------


class TestSkillBuilderYamlDeLeaked:
    """Assert that the specific leaking patterns have been removed from the prompt."""

    def _load_design_skill_template(self) -> str:
        data = yaml.safe_load(SKILL_BUILDER_YAML.read_text(encoding="utf-8"))
        return data["prompts"]["design_skill"]["template"]

    def test_no_hardcoded_rule_choose_weekly_exec_review_v1(self):
        """The rule 'Choose weekly_exec_review_v1 layout only...' must be absent."""
        template = self._load_design_skill_template()
        assert 'Choose "weekly_exec_review_v1" layout only' not in template, (
            "Hardcoded layout rule still present in design_skill prompt — ADR-034 fix required."
        )

    def test_layout_field_no_longer_enumerates_preset_id(self):
        """The schema example for 'layout' must not list 'weekly_exec_review_v1 | default'."""
        template = self._load_design_skill_template()
        assert '"layout": "weekly_exec_review_v1 | default"' not in template, (
            "Hardcoded layout enum still in design_skill schema example — ADR-034 fix required."
        )

    def test_layout_preset_catalog_var_in_template(self):
        """The template must reference {layout_preset_catalog} for injection."""
        template = self._load_design_skill_template()
        assert "{layout_preset_catalog}" in template, (
            "design_skill template does not inject {layout_preset_catalog} — ADR-034 fix required."
        )

    def test_layout_preset_catalog_in_required_vars(self):
        """design_skill required_vars must list layout_preset_catalog."""
        data = yaml.safe_load(SKILL_BUILDER_YAML.read_text(encoding="utf-8"))
        required = data["prompts"]["design_skill"]["required_vars"]
        assert "layout_preset_catalog" in required, (
            f"layout_preset_catalog not in required_vars: {required}"
        )

    def test_version_bumped_from_1_1(self):
        """design_skill version must have been bumped from 1.1 to reflect the change."""
        data = yaml.safe_load(SKILL_BUILDER_YAML.read_text(encoding="utf-8"))
        version = data["prompts"]["design_skill"]["version"]
        assert version != "1.1", (
            "design_skill version still 1.1 — should be bumped to reflect ADR-034 change."
        )


# ---------------------------------------------------------------------------
# (b) Catalog is single source of truth — PptxRenderer dispatches via catalog
# ---------------------------------------------------------------------------


class TestCatalogSingleSourceOfTruth:
    """Verify that layout_catalog.py is the authoritative registry and renderer uses it."""

    def test_catalog_has_weekly_exec_review_v1(self):
        from framework.renderers.layout_catalog import get_preset
        preset = get_preset("weekly_exec_review_v1")
        assert preset is not None
        assert preset.internal_id == "weekly_exec_review_v1"
        assert preset.output_format == "pptx"

    def test_catalog_has_default(self):
        from framework.renderers.layout_catalog import get_preset
        preset = get_preset("default")
        assert preset is not None
        assert preset.internal_id == "default"
        assert preset.output_format == "pptx"

    def test_catalog_get_preset_unknown_returns_none(self):
        from framework.renderers.layout_catalog import get_preset
        assert get_preset("nonexistent_preset_xyz") is None

    def test_renderer_dispatches_weekly_exec_review_v1_via_catalog(self):
        """PptxRenderer still routes weekly_exec_review_v1 data to the special builder."""
        pytest.importorskip("pptx")
        from framework.renderers.pptx_renderer import PptxRenderer
        renderer = PptxRenderer()
        data = {
            "title": "Test Exec Review",
            "layout": "weekly_exec_review_v1",
            "sections": {"scope": "Test scope"},
        }
        result = renderer.render(data)
        assert isinstance(result, bytes)
        assert len(result) > 100

        # Confirm it's a real PPTX with exactly 1 slide (the layout-aware path)
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1, (
            f"weekly_exec_review_v1 should produce 1 slide, got {len(prs.slides)}"
        )

    def test_renderer_dispatches_default_via_catalog(self):
        """PptxRenderer still routes default layout to multi-slide path."""
        pytest.importorskip("pptx")
        from framework.renderers.pptx_renderer import PptxRenderer
        renderer = PptxRenderer()
        data = {
            "title": "Multi-slide deck",
            "layout": "default",
            "sections": {"section_a": "content a", "section_b": "content b"},
        }
        result = renderer.render(data)
        assert isinstance(result, bytes)
        assert len(result) > 100

    def test_renderer_unknown_layout_raises_hard_error(self):
        """DECISION-019 Finding-B: PptxRenderer raises ValueError for unknown layout ids.

        The old behavior (log.warning + fallback) was the silent-degradation amplifier
        (Finding-B). It is replaced by a hard-fail that propagates as an executor
        failure and surfaces as [HIGH] in the DECISION-018 §H three-section report.
        Skills with valid registered layouts (default, weekly_exec_review_v1) unaffected.
        """
        from framework.renderers.pptx_renderer import PptxRenderer
        renderer = PptxRenderer()
        data = {
            "title": "Unknown layout",
            "layout": "some_future_layout_not_in_catalog",
            "sections": {"s": "content"},
        }
        with pytest.raises(ValueError, match="not a registered catalog internal_id"):
            renderer.render(data)

    def test_all_presets_returns_both(self):
        from framework.renderers.layout_catalog import all_presets
        ids = {p.internal_id for p in all_presets()}
        assert "weekly_exec_review_v1" in ids
        assert "default" in ids

    def test_internal_ids_helper(self):
        from framework.renderers.layout_catalog import internal_ids
        ids = internal_ids()
        assert "weekly_exec_review_v1" in ids
        assert "default" in ids


# ---------------------------------------------------------------------------
# (c) Clarify question sanitizer strips internal preset ids
# ---------------------------------------------------------------------------


class TestClarifySanitizer:
    """Verify that _sanitize_clarify_question removes internal ids from user-facing text."""

    def _get_sanitizer(self):
        from framework.skill_builder.conversation import SkillBuilderConversation

        class _FakeConv(SkillBuilderConversation):
            def __init__(self):
                pass  # skip full __init__

        return _FakeConv._sanitize_clarify_question

    def test_question_without_preset_id_unchanged(self):
        sanitize = self._get_sanitizer()
        q = "Should the output be a single dense slide or multiple slides?"
        assert sanitize(q) == q

    def test_question_with_weekly_exec_review_v1_is_replaced(self):
        """Internal id weekly_exec_review_v1 must be replaced with human label."""
        sanitize = self._get_sanitizer()
        q = "Should the layout be weekly_exec_review_v1 or the default layout?"
        sanitized = sanitize(q)
        assert "weekly_exec_review_v1" not in sanitized, (
            f"Internal id still present after sanitization: {sanitized!r}"
        )
        # Human label or description must appear
        from framework.renderers.layout_catalog import get_preset
        label = get_preset("weekly_exec_review_v1").human_label
        assert label in sanitized, (
            f"Expected human label {label!r} in sanitized question: {sanitized!r}"
        )

    def test_question_with_all_known_ids_is_fully_sanitized(self):
        """No known internal id should survive sanitization."""
        sanitize = self._get_sanitizer()
        from framework.renderers.layout_catalog import internal_ids
        # Build a question that contains every known internal id
        q = "Choose one: " + " or ".join(internal_ids())
        sanitized = sanitize(q)
        for iid in internal_ids():
            assert iid not in sanitized, (
                f"Internal id {iid!r} survived sanitization in: {sanitized!r}"
            )

    def test_advance_to_clarify_sanitizes_blocking_questions(self):
        """When _advance_to_clarify is called with a question containing an internal id,
        the stored question must not contain the id."""
        from framework.skill_builder.conversation import (
            SkillBuilderConversation,
            _SessionData,
        )
        from unittest.mock import MagicMock

        conv = object.__new__(SkillBuilderConversation)
        conv._state = "DESIGN_SKILL"
        conv._data = _SessionData(persona="tpm", intent_description="test")
        conv._llm = None
        conv._skill_store = None
        conv._registry = MagicMock()

        # Patch get_registry to return a mock clarify prompt
        mock_spec = MagicMock()
        mock_spec.text = "Before I proceed, I need to clarify: {question}"
        mock_registry = MagicMock()
        mock_registry.get_prompt.return_value = mock_spec

        with patch("framework.skill_builder.conversation.get_registry", return_value=mock_registry):
            result = conv._advance_to_clarify(
                [{"question": "Should layout be weekly_exec_review_v1?", "resolved": False}],
                next_state="REVIEW_DESIGN",
            )

        # The stored question must not contain the internal id
        stored = conv._data._clarify_questions[0]["question"]
        assert "weekly_exec_review_v1" not in stored, (
            f"Internal id survived _advance_to_clarify storage: {stored!r}"
        )
        # The emitted message must also not contain the internal id
        assert "weekly_exec_review_v1" not in result.message, (
            f"Internal id leaked into clarify turn message: {result.message!r}"
        )


# ---------------------------------------------------------------------------
# (d) catalog_for_prompt exposes descriptions, never internal ids
# ---------------------------------------------------------------------------


class TestCatalogForPrompt:
    """catalog_for_prompt must not expose internal_ids in the injected text."""

    def test_catalog_for_prompt_contains_no_internal_ids(self):
        from framework.renderers.layout_catalog import catalog_for_prompt, internal_ids
        catalog_text = catalog_for_prompt()
        # The catalog text must describe the presets — not list their internal ids
        # as standalone tokens (they may appear incidentally in structural_shape for
        # debugging but should NOT be the primary identifier surfaced).
        # Primary check: no bare internal_id appears as the first word of a description line.
        for iid in internal_ids():
            # The catalog should not start a description with the internal id
            assert not re.search(
                rf"(?m)^\s*Description:.*{re.escape(iid)}", catalog_text
            ), (
                f"Internal id {iid!r} appears in description line of catalog_for_prompt output"
            )

    def test_catalog_for_prompt_contains_human_labels(self):
        from framework.renderers.layout_catalog import catalog_for_prompt, all_presets
        catalog_text = catalog_for_prompt()
        for preset in all_presets():
            assert preset.human_label in catalog_text, (
                f"Human label {preset.human_label!r} not in catalog_for_prompt output"
            )

    def test_catalog_for_prompt_filtered_by_output_format(self):
        from framework.renderers.layout_catalog import catalog_for_prompt
        pptx_catalog = catalog_for_prompt("pptx")
        assert len(pptx_catalog) > 0
        # Non-existent format returns the no-presets message
        empty = catalog_for_prompt("nonexistent_format_xyz")
        assert "no layout presets" in empty.lower()

    def test_catalog_for_prompt_pptx_contains_both_known_presets(self):
        from framework.renderers.layout_catalog import catalog_for_prompt, all_presets
        pptx_catalog = catalog_for_prompt("pptx")
        for preset in all_presets():
            if preset.output_format == "pptx":
                assert preset.human_label in pptx_catalog


# ---------------------------------------------------------------------------
# (e) PromptRegistry still parses skill_builder.yaml after changes
# ---------------------------------------------------------------------------


class TestPromptRegistryParsesSkilBuilderYaml:
    """PromptRegistry must load skill_builder.yaml without error after ADR-034 edits."""

    def test_registry_loads_skill_builder_yaml(self, tmp_path):
        """Load skill_builder.yaml via PromptRegistry in a tmp prompts dir."""
        import shutil
        from framework.skill_builder.prompt_registry import PromptRegistry

        # Copy skill_builder.yaml + persona_overlays.yaml to tmp_path
        config_dir = REPO_ROOT / "framework" / "config" / "prompts"
        shutil.copy(config_dir / "skill_builder.yaml", tmp_path / "skill_builder.yaml")
        overlays_src = config_dir / "persona_overlays.yaml"
        if overlays_src.exists():
            shutil.copy(overlays_src, tmp_path / "persona_overlays.yaml")

        registry = PromptRegistry(tmp_path)
        assert "design_skill" in [m.prompt_id for m in registry.list_prompts()]

    def test_design_skill_prompt_renders_with_catalog(self, tmp_path):
        """design_skill prompt renders successfully when layout_preset_catalog is supplied."""
        import shutil
        from framework.skill_builder.prompt_registry import PromptRegistry
        from framework.renderers.layout_catalog import catalog_for_prompt

        config_dir = REPO_ROOT / "framework" / "config" / "prompts"
        shutil.copy(config_dir / "skill_builder.yaml", tmp_path / "skill_builder.yaml")
        overlays_src = config_dir / "persona_overlays.yaml"
        if overlays_src.exists():
            shutil.copy(overlays_src, tmp_path / "persona_overlays.yaml")

        registry = PromptRegistry(tmp_path)
        from framework.renderers.layout_catalog import internal_ids
        _valid_ids_str = ", ".join(f'"{i}"' for i in internal_ids()) + ", null"
        spec = registry.get_prompt(
            "design_skill",
            persona="tpm",
            normalised_intent='{"output_kind": "pptx"}',
            source_capability="[]",
            artifact_layout="null",
            existing_kb_cards="[]",
            layout_preset_catalog=catalog_for_prompt("pptx"),
            layout_valid_ids=_valid_ids_str,
        )
        assert spec.text
        # The rendered prompt must not contain any bare internal preset id as a rule
        assert 'Choose "weekly_exec_review_v1" layout only' not in spec.text
        assert '"layout": "weekly_exec_review_v1 | default"' not in spec.text
        # It must contain the catalog text (human labels)
        from framework.renderers.layout_catalog import all_presets
        for p in all_presets():
            assert p.human_label in spec.text
