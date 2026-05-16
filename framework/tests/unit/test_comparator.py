"""Unit tests for ArtifactComparator (P2 — ADR-029 implementation plan).

Coverage:
    test_is_image_only_true         — image-only PPTX (no text shapes) → True
    test_is_image_only_false        — PPTX with text → False
    test_is_image_only_md_empty     — empty MD bytes → True
    test_is_image_only_md_nonempty  — MD with content → False
    test_is_image_only_unsupported  — unknown type raises ValueError
    test_structure_score_perfect    — identical sections → structure_score == 1.0
    test_structure_score_missing_sections — produced missing 3 of 7 → score == 4/7
    test_density_score_thin_section — produced section has 20 words vs 100 → thin
    test_synonym_normalisation      — "Next Steps" ≈ "Action Items" counts as match
    test_unsupported_type_compare   — unsupported type raises ValueError
    test_gap_report_contains_missing — gap_report names missing sections
    test_gap_report_perfect_match   — gap_report notes no structural gaps
    test_comparator_result_to_dict  — ComparatorResult.to_dict() round-trips
    test_image_only_signal_type     — IMAGE_ONLY_MESSAGE is exported (S5 uses it)
"""
from __future__ import annotations

import io
from unittest.mock import MagicMock

import pytest

from framework.skill_builder.comparator import (
    ArtifactComparator,
    ComparatorResult,
    IMAGE_ONLY_MESSAGE,
    _normalise_section,
)


# ---------------------------------------------------------------------------
# PPTX fixture builders — mirrors the pattern in test_adr026_source_grounded_review.py
# ---------------------------------------------------------------------------

def _make_image_only_pptx_bytes() -> bytes:
    """Build a minimal PPTX with a single rectangle shape (no text) → image-only."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank slide
    slide = prs.slides.add_slide(blank_layout)
    # Add a shape with no text (MSO_SHAPE_TYPE 1 = rectangle-ish auto-shape)
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(4), Inches(3))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_text_pptx_bytes(titles: list[str], body_words_per_slide: int = 50) -> bytes:
    """Build a PPTX with titled slides and body text.

    Args:
        titles:               list of slide titles (one slide per title).
        body_words_per_slide: number of filler words in each slide body.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    content_layout = prs.slide_layouts[1]  # Title and Content

    for title in titles:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        # Fill body placeholder with filler words
        body = slide.placeholders[1]
        body.text = " ".join(["word"] * body_words_per_slide)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_md_bytes(content: str) -> bytes:
    return content.encode("utf-8")


# ---------------------------------------------------------------------------
# TestIsImageOnly
# ---------------------------------------------------------------------------

class TestIsImageOnly:
    def test_is_image_only_true(self):
        """Image-only PPTX (no text shapes) → True."""
        comparator = ArtifactComparator()
        pptx_bytes = _make_image_only_pptx_bytes()
        assert comparator.is_image_only(pptx_bytes, "pptx") is True

    def test_is_image_only_false(self):
        """PPTX with text shapes → False."""
        comparator = ArtifactComparator()
        pptx_bytes = _make_text_pptx_bytes(["Project Status", "Risks"])
        assert comparator.is_image_only(pptx_bytes, "pptx") is False

    def test_is_image_only_md_empty(self):
        """Empty MD bytes → True."""
        comparator = ArtifactComparator()
        assert comparator.is_image_only(b"", "md") is True

    def test_is_image_only_md_nonempty(self):
        """MD with real content → False."""
        comparator = ArtifactComparator()
        md = _make_md_bytes("# Status\n\nThis is content.\n")
        assert comparator.is_image_only(md, "md") is False

    def test_is_image_only_txt_empty(self):
        """Empty TXT → True."""
        comparator = ArtifactComparator()
        assert comparator.is_image_only(b"   ", "txt") is True

    def test_is_image_only_txt_nonempty(self):
        """TXT with words → False."""
        comparator = ArtifactComparator()
        assert comparator.is_image_only(b"hello world", "txt") is False

    def test_is_image_only_unsupported_type_raises(self):
        """Unsupported artifact_type raises ValueError."""
        comparator = ArtifactComparator()
        with pytest.raises(ValueError, match="unsupported artifact_type"):
            comparator.is_image_only(b"some bytes", "pdf")

    def test_is_image_only_case_insensitive(self):
        """artifact_type matching is case-insensitive."""
        comparator = ArtifactComparator()
        md = _make_md_bytes("# Title\n\nContent here.\n")
        assert comparator.is_image_only(md, "MD") is False


# ---------------------------------------------------------------------------
# TestStructureScore
# ---------------------------------------------------------------------------

class TestStructureScore:
    def test_structure_score_perfect_match(self):
        """Reference and produced have identical section names → score == 1.0."""
        titles = ["Status", "Risks", "Next Steps", "Timeline"]
        ref_bytes = _make_text_pptx_bytes(titles, body_words_per_slide=50)
        prod_bytes = _make_text_pptx_bytes(titles, body_words_per_slide=50)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.structure_score == 1.0
        assert result.missing_sections == []

    def test_structure_score_missing_sections(self):
        """Produced missing 3 of 7 sections → structure_score == 4/7, missing lists 3 names."""
        ref_titles = ["Status", "Risks", "Next Steps", "Timeline", "Budget", "Team", "Summary"]
        prod_titles = ["Status", "Timeline", "Budget", "Team"]  # missing 3
        ref_bytes = _make_text_pptx_bytes(ref_titles, body_words_per_slide=40)
        prod_bytes = _make_text_pptx_bytes(prod_titles, body_words_per_slide=40)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")

        expected_score = round(4 / 7, 4)
        assert result.structure_score == pytest.approx(expected_score, abs=0.01)
        # The 3 missing sections should all appear in missing_sections
        missing_lower = {s.lower() for s in result.missing_sections}
        assert "risks" in missing_lower
        assert "next steps" in missing_lower
        assert "summary" in missing_lower

    def test_structure_score_all_missing(self):
        """Produced has no matching sections → structure_score == 0.0."""
        ref_bytes = _make_text_pptx_bytes(["Alpha", "Beta", "Gamma"], body_words_per_slide=20)
        prod_bytes = _make_text_pptx_bytes(["Delta", "Epsilon"], body_words_per_slide=20)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.structure_score == 0.0
        assert len(result.missing_sections) == 3


# ---------------------------------------------------------------------------
# TestDensityScore
# ---------------------------------------------------------------------------

class TestDensityScore:
    def test_density_score_thin_section(self):
        """Produced section has 20 words vs reference 100 words → ratio < 0.5 → thin."""
        ref_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=100)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=20)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.density_score < 0.5
        assert "Status" in result.thin_sections

    def test_density_score_adequate_section(self):
        """Produced section has >= 50% of reference words → NOT in thin_sections."""
        ref_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=100)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=60)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.thin_sections == []

    def test_density_score_capped_at_1(self):
        """Produced section denser than reference → density capped at 1.0."""
        ref_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=50)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=200)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        # density_score should be exactly 1.0 (capped) — not > 1.0
        assert result.density_score <= 1.0
        assert result.density_score == pytest.approx(1.0, abs=0.01)


# ---------------------------------------------------------------------------
# TestSynonymNormalisation
# ---------------------------------------------------------------------------

class TestSynonymNormalisation:
    def test_synonym_normalisation_next_steps_action_items(self):
        """'Next Steps' in reference matches 'Action Items' in produced → counted as match."""
        ref_bytes = _make_text_pptx_bytes(["Status", "Next Steps", "Risks"], body_words_per_slide=50)
        prod_bytes = _make_text_pptx_bytes(["Status", "Action Items", "Risks"], body_words_per_slide=50)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        # All 3 sections should be matched (via synonym map)
        assert result.structure_score == pytest.approx(1.0, abs=0.01)
        assert result.missing_sections == []

    def test_synonym_normalisation_milestones(self):
        """'Key Milestones' in reference matches 'Timeline' in produced."""
        ref_bytes = _make_text_pptx_bytes(["Key Milestones"], body_words_per_slide=30)
        prod_bytes = _make_text_pptx_bytes(["Timeline"], body_words_per_slide=30)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.structure_score == pytest.approx(1.0, abs=0.01)

    def test_synonym_normalisation_risks_and_mitigations(self):
        """'Risks & Mitigations' normalises to match 'Risks'."""
        ref_bytes = _make_text_pptx_bytes(["Risks & Mitigations"], body_words_per_slide=40)
        prod_bytes = _make_text_pptx_bytes(["Risks"], body_words_per_slide=40)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.structure_score == pytest.approx(1.0, abs=0.01)

    def test_non_synonym_still_missing(self):
        """Different non-synonym names are correctly flagged as missing."""
        ref_bytes = _make_text_pptx_bytes(["Governance"], body_words_per_slide=30)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=30)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert result.structure_score == 0.0
        assert "Governance" in result.missing_sections

    def test_normalise_section_function(self):
        """Unit-test _normalise_section directly for synonym resolution."""
        assert _normalise_section("Next Steps") == _normalise_section("Action Items")
        assert _normalise_section("Key Milestones") == _normalise_section("Timeline")
        assert _normalise_section("Risks & Mitigations") == _normalise_section("Risks")
        # Unknown term normalises to itself (lowercased, spaces collapsed)
        assert _normalise_section("Governance") == "governance"


# ---------------------------------------------------------------------------
# TestMarkdownComparison
# ---------------------------------------------------------------------------

class TestMarkdownComparison:
    def test_md_structure_score(self):
        """Markdown section comparison respects heading hierarchy."""
        ref_md = _make_md_bytes(
            "# Status\n\nContent here.\n\n"
            "# Risks\n\nRisk content.\n\n"
            "# Next Steps\n\nStep content.\n"
        )
        prod_md = _make_md_bytes(
            "# Status\n\nProduced content.\n\n"
            "# Next Steps\n\nProduced steps.\n"
            # Risks missing
        )
        comparator = ArtifactComparator()
        result = comparator.compare(ref_md, prod_md, "md")
        assert result.structure_score == pytest.approx(2 / 3, abs=0.02)
        assert "Risks" in result.missing_sections

    def test_txt_treated_as_md(self):
        """TXT artifact type uses the markdown extractor (heading-based)."""
        ref_txt = _make_md_bytes("# Alpha\n\nSome content.\n\n# Beta\n\nMore content.\n")
        prod_txt = _make_md_bytes("# Alpha\n\nProduced content.\n")
        comparator = ArtifactComparator()
        result = comparator.compare(ref_txt, prod_txt, "txt")
        assert result.structure_score < 1.0
        assert "Beta" in result.missing_sections


# ---------------------------------------------------------------------------
# TestUnsupportedType
# ---------------------------------------------------------------------------

class TestUnsupportedType:
    def test_compare_unsupported_type_raises(self):
        """Unsupported artifact_type raises ValueError on compare()."""
        comparator = ArtifactComparator()
        with pytest.raises(ValueError, match="unsupported artifact_type"):
            comparator.compare(b"ref", b"prod", "pdf")

    def test_compare_unsupported_type_is_actionable(self):
        """ValueError message names the supported types."""
        comparator = ArtifactComparator()
        with pytest.raises(ValueError) as exc_info:
            comparator.compare(b"ref", b"prod", "xlsx")
        assert "pptx" in str(exc_info.value)


# ---------------------------------------------------------------------------
# TestGapReport
# ---------------------------------------------------------------------------

class TestGapReport:
    def test_gap_report_contains_missing_sections(self):
        """gap_report names the missing sections by name."""
        ref_bytes = _make_text_pptx_bytes(["Status", "Risks", "Next Steps"], body_words_per_slide=50)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=50)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert "Risks" in result.gap_report
        assert "Next Steps" in result.gap_report
        assert "Missing" in result.gap_report

    def test_gap_report_perfect_match_notes_no_gaps(self):
        """gap_report notes no structural gaps when all sections match."""
        titles = ["Status", "Risks", "Next Steps"]
        ref_bytes = _make_text_pptx_bytes(titles, body_words_per_slide=50)
        prod_bytes = _make_text_pptx_bytes(titles, body_words_per_slide=60)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert "No structural gaps" in result.gap_report

    def test_gap_report_mentions_thin_sections(self):
        """gap_report flags thin sections when density < 0.5."""
        ref_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=100)
        prod_bytes = _make_text_pptx_bytes(["Status"], body_words_per_slide=20)
        comparator = ArtifactComparator()
        result = comparator.compare(ref_bytes, prod_bytes, "pptx")
        assert "Thin" in result.gap_report or "thin" in result.gap_report


# ---------------------------------------------------------------------------
# TestComparatorResultDataclass
# ---------------------------------------------------------------------------

class TestComparatorResultDataclass:
    def test_comparator_result_to_dict(self):
        """ComparatorResult.to_dict() exposes all required keys for EVAL turn."""
        result = ComparatorResult(
            structure_score=0.75,
            density_score=0.5,
            missing_sections=["Risks"],
            thin_sections=["Status"],
            gap_report="The produced PPTX has 3 sections; your reference had 4.",
        )
        d = result.to_dict()
        assert d["structure_score"] == 0.75
        assert d["density_score"] == 0.5
        assert d["missing_sections"] == ["Risks"]
        assert d["thin_sections"] == ["Status"]
        assert "gap_report" in d

    def test_comparator_result_defaults(self):
        """ComparatorResult initialises with empty lists by default."""
        result = ComparatorResult(structure_score=1.0, density_score=1.0)
        assert result.missing_sections == []
        assert result.thin_sections == []
        assert result.gap_report == ""


# ---------------------------------------------------------------------------
# TestImageOnlySignal
# ---------------------------------------------------------------------------

class TestImageOnlySignal:
    def test_image_only_message_exported(self):
        """IMAGE_ONLY_MESSAGE is importable and contains the required user text."""
        assert "Image-based reference artifacts" in IMAGE_ONLY_MESSAGE
        assert "Vision-LLM" in IMAGE_ONLY_MESSAGE
        assert "text-bearing" in IMAGE_ONLY_MESSAGE

    def test_image_only_message_matches_adr029_prescription(self):
        """IMAGE_ONLY_MESSAGE matches the ADR-029 §C.1 prescribed message exactly."""
        # ADR-029 prescribes this exact wording; S5 surfaces it verbatim.
        assert "no Vision-LLM backend" in IMAGE_ONLY_MESSAGE
        assert "Please upload a text-bearing reference" in IMAGE_ONLY_MESSAGE


# ---------------------------------------------------------------------------
# TestEmptyReferenceError
# ---------------------------------------------------------------------------

class TestEmptyReferenceError:
    def test_empty_reference_raises(self):
        """compare() with a zero-section reference raises ValueError, not silently 1.0."""
        # MD with no headings → no sections
        ref_bytes = _make_md_bytes("just plain text no headings\n")
        prod_bytes = _make_md_bytes("# Status\n\nContent.\n")
        comparator = ArtifactComparator()
        # No headings → _extract_sections returns [] → ValueError
        with pytest.raises(ValueError, match="no extractable sections"):
            comparator.compare(ref_bytes, prod_bytes, "md")
