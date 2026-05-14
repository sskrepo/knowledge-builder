"""PPTX renderer — produces a real PowerPoint deck via python-pptx.

Per ADR-016. We have working precedent for python-pptx in this session
(used for the executive brief). This wraps it in the Renderer Protocol.

ADR-026 Fix 5: layout-aware rendering.
When data["layout"] == "weekly_exec_review_v1", delegates to
_render_weekly_exec_review_v1() which builds the single-slide two-column
Oracle-style layout programmatically (no binary template dependency).

All other layouts follow the original multi-slide path.
"""
from __future__ import annotations
import io
import logging
import re
from pathlib import Path
from ._base import BaseRenderer

log = logging.getLogger(__name__)

# Oracle brand palette (from Oracle corporate guidelines)
_ORACLE_RED = (199, 70, 52)       # #C74634
_ORACLE_DARK = (61, 61, 61)       # #3D3D3D
_ORACLE_LIGHT_BG = (248, 248, 248)  # #F8F8F8 — sidebar background
_ORACLE_BORDER = (220, 220, 220)  # #DCDCDC — table borders

# Status keyword → Oracle-ish bold (we bold the keyword inline)
_STATUS_KEYWORDS = [
    "Completed", "Complete", "Done",
    "Approved", "Approved",
    "In Progress", "In-Progress",
    "On Hold", "Hold",
    "At Risk", "At-Risk",
    "Blocked",
    "Slipped",
    "Green", "Amber", "Red",
]


class PptxRenderer(BaseRenderer):
    name = "pptx"
    output_extension = ".pptx"

    def render(self, data: dict, template: str | Path | None = None) -> bytes:
        layout = data.get("layout", "")
        if layout == "weekly_exec_review_v1":
            return self._render_weekly_exec_review_v1(data)
        return self._render_default(data, template)

    # ------------------------------------------------------------------
    # Layout: weekly_exec_review_v1 (ADR-026 Fix 5)
    # ------------------------------------------------------------------

    def _render_weekly_exec_review_v1(self, data: dict) -> bytes:
        """Build the single-slide two-column Oracle-style weekly exec review.

        Slide layout (widescreen 13.33" x 7.5"):
          - Title top-left (bold, large)
          - Jira ID top-right (accent, smaller)
          - Left column (0.3" → 7.2" wide, 1.1" → 5.9" tall):
              2-row table:
                Row 1 — Scope (thin row)
                Row 2 — Assumptions + Status bullets + Next Steps
          - Right sidebar (7.6" → 5.4" wide, 1.1" → 5.9" tall):
              3 stacked boxes: Key Milestones, ORM Status, Risk/Mitigation
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            log.warning("python-pptx not installed; producing placeholder bytes")
            return self._stub_bytes(data)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Use blank slide layout (index 6) — no placeholder interference
        blank_layout = None
        for layout in prs.slide_layouts:
            if layout.name in ("Blank", "blank"):
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        slide = prs.slides.add_slide(blank_layout)

        # --- Oracle header band (top, full width, thin red bar) ---
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        header_shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.45),
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(*_ORACLE_RED)
        header_shape.line.fill.background()  # no border

        # --- Title ---
        title_text = data.get("title") or data.get("sections", {}).get("Project Name", "Weekly Exec Review")
        tx_title = slide.shapes.add_textbox(
            Inches(0.25), Inches(0.5),
            Inches(9.0), Inches(0.55),
        )
        tf = tx_title.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(title_text)
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(*_ORACLE_DARK)

        # --- Jira ID top-right ---
        jira_id = (
            data.get("jira_id")
            or data.get("sections", {}).get("Jira Id")
            or data.get("extracted", {}).get("jira_id", "")
        )
        if jira_id:
            tx_jira = slide.shapes.add_textbox(
                Inches(9.4), Inches(0.5),
                Inches(3.7), Inches(0.5),
            )
            tf_j = tx_jira.text_frame
            tf_j.word_wrap = False
            pj = tf_j.paragraphs[0]
            pj.alignment = PP_ALIGN.RIGHT
            rj = pj.add_run()
            rj.text = str(jira_id)
            rj.font.bold = True
            rj.font.size = Pt(11)
            rj.font.color.rgb = RGBColor(*_ORACLE_RED)

        # ------------------------------------------------------------------
        # Left column: 2-row table
        # ------------------------------------------------------------------
        # Extract fields from either data["extracted"] or data["sections"].
        # _get and _get_list accept multiple keys/aliases so the layout works
        # with both the ADR-026 canonical schema field names AND the
        # natural-source names that ADR-027 DESIGN_SKILL tends to choose
        # (e.g. 'milestones' instead of 'key_milestones', 'in_scope' instead
        # of 'scope'). The first non-empty value wins.
        def _resolve_keys(keys: tuple[str, ...]) -> object:
            extracted = data.get("extracted", {}) or {}
            sections = data.get("sections", {}) or {}
            for k in keys:
                v = extracted.get(k)
                if v not in (None, "", []):
                    return v
                v2 = sections.get(k.replace("_", " ").title())
                if v2 not in (None, "", []):
                    return v2
            return None

        def _get(*keys: str) -> str:
            v = _resolve_keys(keys)
            return _format_value(v) if v is not None else ""

        def _get_list(*keys: str) -> list[str]:
            v = _resolve_keys(keys)
            return _to_bullets(v) if v is not None else []

        # Scope: prefer explicit scope; otherwise compose from in_scope +
        # out_of_scope (the way the source page often expresses it).
        scope_text = _get("scope", "in_scope", "business_outcome")
        out_of_scope = _get("out_of_scope")
        if scope_text and out_of_scope:
            scope_text = f"In scope: {scope_text}. Out of scope: {out_of_scope}."
        assumptions     = _get_list("assumptions")
        status_bullets  = _get_list("status_bullets", "weekly_status_update",
                                    "current_status", "current_phase", "overall_status")
        next_steps      = _get_list("next_steps", "completed_last_week", "next_week_plan")
        key_milestones  = _get_list("key_milestones", "milestones", "top_milestones")
        orm_status      = _get("orm_status", "orm")
        risks           = _get_list("risks_mitigations", "risks", "top_risks",
                                    "top_issues_blockers", "blockers")

        # Left column dimensions
        lc_left = Inches(0.25)
        lc_top = Inches(1.1)
        lc_width = Inches(7.2)
        lc_height = Inches(5.9)

        # Table: 2 rows × 1 column
        from pptx.util import Inches
        tbl = slide.shapes.add_table(2, 1, lc_left, lc_top, lc_width, lc_height).table

        # Row 0 height: ~0.8" for scope
        from pptx.util import Inches as _I
        tbl.rows[0].height = _I(0.75)
        # Row 1 height: remaining space (python-pptx auto-sizes)

        # Style table border and fills
        for row_idx in range(2):
            cell = tbl.cell(row_idx, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Row 0 — Scope
        _add_section_cell(tbl.cell(0, 0), "Scope", scope_text or "(scope not extracted)")

        # Row 1 — Assumptions + Status + Next Steps
        row1_tf = tbl.cell(1, 0).text_frame
        row1_tf.word_wrap = True
        row1_tf.clear()
        p_idx = 0

        def _add_section_header(tf, label: str, para_idx: int) -> int:
            p = tf.paragraphs[para_idx] if para_idx == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = label
            r.font.bold = True
            r.font.size = Pt(10)
            r.font.color.rgb = RGBColor(*_ORACLE_RED)
            return para_idx + 1

        def _add_bullets(tf, items: list[str], para_idx: int) -> int:
            for item in items:
                p = tf.paragraphs[para_idx] if para_idx == 0 else tf.add_paragraph()
                p.space_before = Pt(2)
                _add_bullet_run(p, item)
                para_idx += 1
            return para_idx

        if assumptions:
            p_idx = _add_section_header(row1_tf, "Assumptions", p_idx)
            p_idx = _add_bullets(row1_tf, assumptions, p_idx)
        if status_bullets:
            p_idx = _add_section_header(row1_tf, "Status", p_idx)
            p_idx = _add_bullets(row1_tf, status_bullets, p_idx)
        if next_steps:
            p_idx = _add_section_header(row1_tf, "Next Steps", p_idx)
            _add_bullets(row1_tf, next_steps, p_idx)

        # ------------------------------------------------------------------
        # Right sidebar: 3 stacked boxes
        # ------------------------------------------------------------------
        rs_left = Inches(7.7)
        rs_top = Inches(1.1)
        rs_width = Inches(5.3)
        box_heights = [Inches(2.6), Inches(0.9), Inches(2.4)]  # milestones, orm, risk
        current_top = rs_top

        # Box 1 — Key Milestones
        _add_sidebar_box(slide, rs_left, current_top, rs_width, box_heights[0],
                         "Key Milestones", key_milestones)
        current_top += box_heights[0] + Inches(0.1)

        # Box 2 — ORM Status
        _add_sidebar_box(slide, rs_left, current_top, rs_width, box_heights[1],
                         "ORM", [orm_status] if orm_status else ["(not specified)"])
        current_top += box_heights[1] + Inches(0.1)

        # Box 3 — Risk / Mitigation
        _add_sidebar_box(slide, rs_left, current_top, rs_width, box_heights[2],
                         "Risk / Mitigation", risks)

        # --- Footer: generation note ---
        generated_at = data.get("generated_at", "")
        citations = data.get("citations", [])
        footer_parts = []
        if citations:
            footer_parts.append(f"Source: {citations[0]}")
        if generated_at:
            footer_parts.append(f"Generated {generated_at[:10]}")
        if footer_parts:
            tx_footer = slide.shapes.add_textbox(
                Inches(0.25), Inches(7.1),
                Inches(13.0), Inches(0.3),
            )
            tf_footer = tx_footer.text_frame
            pf = tf_footer.paragraphs[0]
            rf = pf.add_run()
            rf.text = "  |  ".join(footer_parts)
            rf.font.size = Pt(7)
            rf.font.color.rgb = RGBColor(150, 150, 150)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Default multi-slide renderer (original behavior)
    # ------------------------------------------------------------------

    def _render_default(self, data: dict, template: str | Path | None = None) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            log.warning("python-pptx not installed; producing placeholder bytes")
            return self._stub_bytes(data)

        # Start from template if given, else blank widescreen
        if template and Path(template).exists():
            prs = Presentation(str(template))
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = data.get("title", "Generated by Knowledge Builder Framework")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = data.get("subtitle", "")

        # Content slides — one per section
        for section_name, body in data.get("sections", {}).items():
            content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title:
                slide.shapes.title.text = section_name
            # Body
            tf = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx in (1, 2):
                    tf = ph.text_frame
                    break
            if tf is None:
                continue
            tf.clear()
            if isinstance(body, list):
                for i, item in enumerate(body):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"• {item}"
            elif isinstance(body, dict):
                first = True
                for k, v in body.items():
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    p.text = f"• {k}: {v}"
                    first = False
            else:
                tf.text = str(body)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @staticmethod
    def _stub_bytes(data: dict) -> bytes:
        # When python-pptx isn't installed, emit a tiny "manifest" so downstream
        # delivery can still happen — useful for shape-testing the pipeline.
        import json
        return ("STUB PPTX (install python-pptx to render real)\n" +
                json.dumps(data, indent=2, default=str)).encode("utf-8")


# ------------------------------------------------------------------
# Module-level helpers for layout builders
# ------------------------------------------------------------------

def _format_value(val) -> str:
    if isinstance(val, list):
        return "\n".join(f"• {v}" for v in val)
    return str(val)


def _to_bullets(val) -> list[str]:
    if isinstance(val, list):
        return [str(v) for v in val]
    if isinstance(val, str):
        # Split on newlines / bullet chars
        lines = re.split(r"[\n•\-\*]", val)
        return [ln.strip() for ln in lines if ln.strip()]
    return [str(val)]


def _add_section_cell(cell, header: str, body: str) -> None:
    """Populate a table cell with a header + body text."""
    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    r_h = p0.add_run()
    r_h.text = header
    r_h.font.bold = True
    r_h.font.size = Pt(10)
    r_h.font.color.rgb = RGBColor(*_ORACLE_RED)
    # Body text in second paragraph
    p1 = tf.add_paragraph()
    r_b = p1.add_run()
    r_b.text = body
    r_b.font.size = Pt(9)
    r_b.font.color.rgb = RGBColor(*_ORACLE_DARK)


def _add_bullet_run(para, text: str) -> None:
    """Add a bullet item to a paragraph, bolding any known status keywords."""
    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return

    # Check for status keywords to bold
    remaining = f"• {text}"
    for kw in _STATUS_KEYWORDS:
        if kw.lower() in text.lower():
            # Find and bold the keyword
            idx = remaining.lower().find(kw.lower())
            if idx >= 0:
                prefix = remaining[:idx]
                keyword = remaining[idx:idx + len(kw)]
                suffix = remaining[idx + len(kw):]
                if prefix:
                    r = para.add_run()
                    r.text = prefix
                    r.font.size = Pt(9)
                    r.font.color.rgb = RGBColor(*_ORACLE_DARK)
                r_kw = para.add_run()
                r_kw.text = keyword
                r_kw.font.bold = True
                r_kw.font.size = Pt(9)
                r_kw.font.color.rgb = RGBColor(*_ORACLE_DARK)
                remaining = suffix
            break

    if remaining:
        r = para.add_run()
        r.text = remaining
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(*_ORACLE_DARK)


def _add_sidebar_box(slide, left, top, width, height, header: str, items: list[str]) -> None:
    """Add a sidebar box with a colored header band and bullet items."""
    try:
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return

    # Header band (thin colored strip above the box)
    hdr_shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        left, top,
        width, Inches(0.28),
    )
    hdr_shape.fill.solid()
    hdr_shape.fill.fore_color.rgb = RGBColor(*_ORACLE_RED)
    hdr_shape.line.fill.background()

    # Header text (white on red)
    tx_hdr = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(0.04),
        width - Inches(0.1), Inches(0.22),
    )
    p_hdr = tx_hdr.text_frame.paragraphs[0]
    r_hdr = p_hdr.add_run()
    r_hdr.text = header
    r_hdr.font.bold = True
    r_hdr.font.size = Pt(9)
    r_hdr.font.color.rgb = RGBColor(255, 255, 255)

    # Body box (light background)
    body_top = top + Inches(0.28)
    body_height = height - Inches(0.28)
    body_shape = slide.shapes.add_shape(
        1,
        left, body_top,
        width, body_height,
    )
    body_shape.fill.solid()
    body_shape.fill.fore_color.rgb = RGBColor(*_ORACLE_LIGHT_BG)
    body_shape.line.color.rgb = RGBColor(*_ORACLE_BORDER)

    # Body text
    tx_body = slide.shapes.add_textbox(
        left + Inches(0.08), body_top + Inches(0.05),
        width - Inches(0.15), body_height - Inches(0.1),
    )
    tx_body.text_frame.word_wrap = True
    tf = tx_body.text_frame
    first = True
    for item in items or ["(none)"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(2)
        _add_bullet_run(p, str(item))
        first = False
